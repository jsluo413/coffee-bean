"""Build the 10-minute presentation deck for 'Beyond the Bean'.

Run: python -m src.build_pptx
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
IMG = DOCS / "img"

ACCENT = RGBColor(0x6F, 0x4E, 0x37)         # coffee brown
TEXT = RGBColor(0x26, 0x27, 0x30)           # near-black
MUTED = RGBColor(0x8A, 0x8A, 0x8A)


def add_title_slide(prs, title, subtitle, author):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = Inches(0.7); top = Inches(2.0)
    tb = slide.shapes.add_textbox(left, top, Inches(12), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(44)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT

    tb = slide.shapes.add_textbox(left, Inches(3.2), Inches(12), Inches(0.8))
    tb.text_frame.text = subtitle
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(22)
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT

    tb = slide.shapes.add_textbox(left, Inches(4.2), Inches(12), Inches(0.5))
    tb.text_frame.text = author
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED


def add_section_heading(slide, heading):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    tf = tb.text_frame
    tf.text = heading
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = ACCENT


def add_bullet_slide(prs, heading, bullets, sub=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_heading(slide, heading)
    if sub:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.5), Inches(0.5))
        tb.text_frame.text = sub
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED
        body_top = Inches(1.7)
    else:
        body_top = Inches(1.3)
    tb = slide.shapes.add_textbox(Inches(0.7), body_top, Inches(12), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + b
        for r in p.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = TEXT
        p.space_after = Pt(8)
    return slide


def add_image_slide(prs, heading, image_path: Path, caption: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_heading(slide, heading)
    # Image fills width with margins; centered.
    img_w = Inches(11.5)
    img_h = Inches(5.6)
    left = Inches((13.33 - 11.5) / 2)
    top = Inches(1.1)
    slide.shapes.add_picture(str(image_path), left, top, width=img_w, height=img_h)
    if caption:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.5))
        tb.text_frame.text = caption
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide


def main():
    findings = json.loads((DOCS / "findings.json").read_text())
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        title="Beyond the Bean",
        subtitle="Global coffee production, quality, and the specialty market",
        author="Jinsheng Luo  ·  Data Visualization Final Project",
    )

    # 2. The question
    add_bullet_slide(
        prs, "The question",
        bullets=[
            "Where coffee is grown is not where it is graded.",
            "Where it is graded is not where it is sold at a premium.",
            "Three datasets, four hypotheses, one dashboard — does the data tell one story or two?",
        ],
    )

    # 3. Data
    add_bullet_slide(
        prs, "Three datasets",
        bullets=[
            "FAO production (1961–2024) — 7,375 country-years (Our World in Data)",
            "CQI cupping database — 1,338 professionally graded lots, 10 sensory dimensions",
            "coffeereview.com — 2,095 specialty consumer reviews (2017–2022)",
            "Linked at the country level via keyword-matching on origin fields.",
        ],
    )

    # 4. H1 — verdict + map
    h1 = findings["h1"]
    sub = (f"Belt mean = {h1['belt_mean_yield']:.2f} t/ha (n={h1['belt_n']})  ·  "
           f"Margin mean = {h1['margin_mean_yield']:.2f} t/ha (n={h1['margin_n']})  ·  "
           f"Welch t = {h1['welch_t']:.2f}, p = {h1['welch_p']:.3f}")
    add_image_slide(
        prs, "H1 · Production is wildly concentrated — but the belt isn't more productive",
        IMG / "h1_choropleth_2022.png",
        caption=sub,
    )

    # 5. H1 — yield vs latitude
    add_image_slide(
        prs, "H1 · Yield varies more within the belt than between belt and margin",
        IMG / "h1_yield_lat.png",
        caption="Hypothesis rejected. Latitude defines where coffee can grow; agronomy defines how much per hectare.",
    )

    # 6. H2 — altitude
    h2 = findings["h2"]
    means = h2["altitude_means"]
    sub = (f"<1000m = {means.get('<1000m', 0):.1f}  ·  >2000m = {means.get('>2000m', 0):.1f}  ·  "
           f"OLS β(altitude std) = {h2['altitude_beta']:+.2f}, p = {h2['altitude_p']:.1e}")
    add_image_slide(
        prs, "H2 · Altitude lifts cupping score — modestly but reliably",
        IMG / "h2_altitude_box.png",
        caption=sub,
    )

    # 7. H2 — regression
    add_image_slide(
        prs, "H2 · Standardized OLS coefficients on Total Cup Points",
        IMG / "h2_regression.png",
        caption=f"R² = {h2['regression_r2']:.2f}, n = {h2['regression_n']}. Altitude is the largest single positive driver; defects, the largest negative.",
    )

    # 8. H3 — price/rating
    h3 = findings["h3"]
    bm = h3["band_means"]
    sub = (f"Mean rating: <$5 = {bm['<$5']:.1f}, $8–12 = {bm['$8-12']:.1f}, "
           f">$20 = {bm['>$20']:.1f}  ·  Spearman ρ = {h3['spearman_corr']:.2f}")
    add_image_slide(
        prs, "H3 · Above ~$12/100g, more money buys very few extra rating points",
        IMG / "h3_price_rating.png",
        caption=sub,
    )

    # 9. H4 — clusters
    add_image_slide(
        prs, "H4 · Flavor language differs systematically by origin",
        IMG / "h4_clusters.png",
        caption="TF-IDF + KMeans clusters of review descriptions, share by top-10 origins. Nuanced support — origin shifts the mix, not the category.",
    )

    # 10. Synthesis
    syn = findings["synthesis"]
    sub = (f"Volume × quality (Spearman) = {syn['corr_volume_quality_spearman']:+.2f}  ·  "
           f"Volume × price = {syn['corr_volume_price_spearman']:+.2f}  ·  "
           f"Quality × price = {syn['corr_quality_price_spearman']:+.2f}")
    add_image_slide(
        prs, "Synthesis · Volume and quality are independent",
        IMG / "synthesis_bubble.png",
        caption=sub,
    )

    # 11. Verdicts table
    add_bullet_slide(
        prs, "Putting the four hypotheses together",
        bullets=[
            "H1 — Belt → higher yield: REJECTED (margin avg is higher)",
            "H2 — Altitude → higher cup score: SUPPORTED (β = +0.44 std, p < 1e-8)",
            "H3 — Diminishing returns above $12/100g: STRONGLY SUPPORTED",
            "H4 — Flavor clusters track origin: SUPPORTED with nuance",
            "Big idea: coffee is two industries — commodity (volume) and specialty (story).",
        ],
    )

    # 12. Limitations
    add_bullet_slide(
        prs, "What the data cannot tell us",
        bullets=[
            "CQI is a submitted-sample database, not a random sample of world coffee.",
            "coffeereview.com is curated specialty — ratings cluster between 88 and 96.",
            "Origin-of-coffee ≠ country-of-roaster; blends resolved to the first listed origin.",
            "All findings are correlational; no causal claims.",
        ],
    )

    # 13. Closing / Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(12), Inches(2))
    tb.text_frame.text = "Thank you — questions?"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(48)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(12), Inches(1))
    tb.text_frame.text = "Live dashboard:  streamlit run streamlit_app.py    ·    Repo: github.com/<you>/coffee-bean"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED

    out_path = DOCS / "presentation.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
